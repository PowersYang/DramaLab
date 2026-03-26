from __future__ import annotations

from datetime import date
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
DOCS_DIR = ROOT / "docs"
IMAGES_DIR = DOCS_DIR / "images"
OUTPUT_PATH = ROOT / "LumenX_项目改造优化汇报_截至2026-03-26.pptx"


BG = RGBColor(246, 248, 252)
PRIMARY = RGBColor(20, 57, 95)
SECONDARY = RGBColor(62, 110, 166)
ACCENT = RGBColor(27, 155, 139)
ACCENT_2 = RGBColor(237, 179, 76)
TEXT = RGBColor(34, 39, 46)
MUTED = RGBColor(92, 103, 115)
WHITE = RGBColor(255, 255, 255)
LIGHT_CARD = RGBColor(255, 255, 255)
DIVIDER = RGBColor(216, 224, 232)


def add_full_background(slide, color: RGBColor) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_header_bar(slide, title: str, subtitle: str | None = None) -> None:
    title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.35), Inches(8.9), Inches(0.55))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = title
    r.font.name = "PingFang SC"
    r.font.size = Pt(24)
    r.font.bold = True
    r.font.color.rgb = PRIMARY

    line = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.6), Inches(0.92), Inches(12.0), Inches(0.03)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT
    line.line.fill.background()

    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(9.0), Inches(0.34), Inches(3.0), Inches(0.35))
        tf = sub_box.text_frame
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.RIGHT
        r = p.add_run()
        r.text = subtitle
        r.font.name = "PingFang SC"
        r.font.size = Pt(10)
        r.font.color.rgb = MUTED


def add_textbox(slide, left, top, width, height, lines, font_size=16, color=TEXT, bold_first=False) -> None:
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(6)
    tf.margin_right = Pt(6)
    tf.margin_top = Pt(4)
    tf.margin_bottom = Pt(4)
    for idx, line in enumerate(lines):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = line
        p.font.name = "PingFang SC"
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.level = 0
        if line.startswith("•"):
            p.bullet = True
        if idx == 0 and bold_first:
            p.font.bold = True


def add_card(slide, left, top, width, height, title, bullets, accent=SECONDARY) -> None:
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = LIGHT_CARD
    shape.line.color.rgb = DIVIDER

    stripe = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, left, top, Inches(0.12), height)
    stripe.fill.solid()
    stripe.fill.fore_color.rgb = accent
    stripe.line.fill.background()

    add_textbox(slide, left + Inches(0.2), top + Inches(0.12), width - Inches(0.3), Inches(0.38), [title], 18, PRIMARY, True)
    add_textbox(slide, left + Inches(0.18), top + Inches(0.52), width - Inches(0.28), height - Inches(0.62), bullets, 12, TEXT)


def fit_image(path: Path, max_w: float, max_h: float) -> tuple[float, float]:
    with Image.open(path) as img:
        w, h = img.size
    ratio = min(max_w / w, max_h / h)
    return w * ratio, h * ratio


def add_image(slide, path: Path, left, top, max_w, max_h) -> None:
    width_px, height_px = fit_image(path, int(max_w * 96), int(max_h * 96))
    slide.shapes.add_picture(str(path), left, top, width=Inches(width_px / 96), height=Inches(height_px / 96))


def add_metric(slide, left, top, width, height, number, label, accent) -> None:
    box = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height)
    box.fill.solid()
    box.fill.fore_color.rgb = WHITE
    box.line.color.rgb = DIVIDER

    top_bar = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, left, top, width, Inches(0.08))
    top_bar.fill.solid()
    top_bar.fill.fore_color.rgb = accent
    top_bar.line.fill.background()

    num_box = slide.shapes.add_textbox(left + Inches(0.12), top + Inches(0.18), width - Inches(0.24), Inches(0.42))
    tf = num_box.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = number
    r.font.name = "Aptos"
    r.font.size = Pt(24)
    r.font.bold = True
    r.font.color.rgb = PRIMARY

    label_box = slide.shapes.add_textbox(left + Inches(0.08), top + Inches(0.62), width - Inches(0.16), Inches(0.35))
    tf = label_box.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = label
    r.font.name = "PingFang SC"
    r.font.size = Pt(10)
    r.font.color.rgb = MUTED


def build_presentation() -> Presentation:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Slide 1
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_full_background(slide, BG)
    banner = IMAGES_DIR / "LumenX Studio Banner.jpeg"
    slide.shapes.add_picture(str(banner), Inches(0), Inches(0), width=Inches(13.333), height=Inches(2.2))

    panel = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.7), Inches(2.35), Inches(11.95), Inches(4.3))
    panel.fill.solid()
    panel.fill.fore_color.rgb = WHITE
    panel.line.color.rgb = DIVIDER

    add_textbox(slide, Inches(1.0), Inches(2.8), Inches(8.5), Inches(0.8), ["LumenX 项目改造优化汇报"], 28, PRIMARY, True)
    add_textbox(
        slide,
        Inches(1.0),
        Inches(3.6),
        Inches(8.7),
        Inches(1.4),
        [
            "汇报范围：截止 2026-03-26 已完成与进行中的本轮改造优化工作",
            "汇报重点：架构重构、基础能力升级、产品链路增强、下一阶段计划",
        ],
        16,
        TEXT,
    )
    tag = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(1.0), Inches(5.2), Inches(2.6), Inches(0.45))
    tag.fill.solid()
    tag.fill.fore_color.rgb = PRIMARY
    tag.line.fill.background()
    add_textbox(slide, Inches(1.08), Inches(5.26), Inches(2.4), Inches(0.2), ["面向领导汇报版"], 12, WHITE)
    add_textbox(slide, Inches(9.7), Inches(5.45), Inches(2.0), Inches(0.5), [f"生成日期：{date(2026, 3, 26).isoformat()}"], 10, MUTED)

    # Slide 2
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_full_background(slide, BG)
    add_header_bar(slide, "一页总览", "截至 2026-03-26")
    add_textbox(
        slide,
        Inches(0.72),
        Inches(1.15),
        Inches(5.6),
        Inches(1.2),
        [
            "本轮工作已从“单点功能补充”升级为“产品链路 + 后端架构 + 交付基础设施”三位一体优化，",
            "既提升了创作流程完整度，也为后续多项目、多租户、数据库化演进打下基础。",
        ],
        16,
        TEXT,
    )
    metrics = [
        ("82", "主干重构累计涉及文件", ACCENT),
        ("7117+", "累计新增代码行", SECONDARY),
        ("10", "新建/收敛服务模块", ACCENT_2),
        ("12", "仓储层文件数", PRIMARY),
        ("4", "当前测试文件", ACCENT),
    ]
    left = 0.72
    for idx, (number, label, accent) in enumerate(metrics):
        add_metric(slide, Inches(left + idx * 2.4), Inches(2.2), Inches(2.1), Inches(1.2), number, label, accent)
    add_card(
        slide,
        Inches(0.72),
        Inches(4.0),
        Inches(3.85),
        Inches(2.5),
        "业务层面的成果",
        [
            "• 打通 Script -> Asset -> Storyboard -> Motion -> Assembly 的端到端体验",
            "• 补齐系列化创作、提示词润色、导出与交互体验优化",
            "• 提高创作者在单平台内完成生产的闭环能力",
        ],
        ACCENT,
    )
    add_card(
        slide,
        Inches(4.74),
        Inches(4.0),
        Inches(3.85),
        Inches(2.5),
        "技术层面的成果",
        [
            "• 形成 API / Service / Workflow / Provider / Repository 分层结构",
            "• 配置读写统一收口，便于开发态与打包态共存",
            "• 引入数据库模型和仓储映射，降低后续扩展成本",
        ],
        SECONDARY,
    )
    add_card(
        slide,
        Inches(8.76),
        Inches(4.0),
        Inches(3.85),
        Inches(2.5),
        "当前在推进的优化",
        [
            "• 从 MySQL 初始化脚本向 PostgreSQL 兼容模型/脚本迁移",
            "• 强化环境配置、启动引导和数据库初始化流程",
            "• 新增仓储与 Schema 相关测试，提升改造可验证性",
        ],
        ACCENT_2,
    )

    # Slide 3
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_full_background(slide, BG)
    add_header_bar(slide, "改造背景与目标")
    add_card(
        slide,
        Inches(0.72),
        Inches(1.2),
        Inches(4.0),
        Inches(4.9),
        "改造前的主要痛点",
        [
            "• 业务逻辑集中在少数大文件，理解成本和修改风险都偏高",
            "• 配置读取、模型调用、存储逻辑存在耦合，难以做平台化扩展",
            "• 面向系列创作、多人协作、长期演进的基础能力不足",
            "• 测试与数据层抽象偏弱，数据库迁移和结构调整成本高",
        ],
        RGBColor(209, 78, 78),
    )
    add_card(
        slide,
        Inches(4.92),
        Inches(1.2),
        Inches(3.55),
        Inches(4.9),
        "本轮改造目标",
        [
            "• 让产品链路更完整，减少创作者在多个工具之间切换",
            "• 让后端结构更清晰，支撑并行开发和后续功能扩展",
            "• 让配置、存储、模型接入更标准，降低环境复杂度",
            "• 让数据库层逐步从脚本化走向规范化、可验证化",
        ],
        ACCENT,
    )
    add_card(
        slide,
        Inches(8.67),
        Inches(1.2),
        Inches(3.95),
        Inches(4.9),
        "本轮汇报关注点",
        [
            "• 已落地：主干架构拆解、系列能力、提示词链路、导出闭环",
            "• 正推进：PostgreSQL 兼容、统一环境配置、仓储测试补强",
            "• 接下来：多租户、组织空间、稳定性和数据迁移继续深化",
        ],
        ACCENT_2,
    )

    # Slide 4
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_full_background(slide, BG)
    add_header_bar(slide, "核心改造一：后端架构重构")

    left_box = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.72), Inches(1.35), Inches(5.7), Inches(5.55))
    left_box.fill.solid()
    left_box.fill.fore_color.rgb = WHITE
    left_box.line.color.rgb = DIVIDER
    add_textbox(slide, Inches(0.98), Inches(1.6), Inches(5.1), Inches(0.45), ["重构方向：从“大而全 service”转向“职责分层”"], 20, PRIMARY, True)
    add_textbox(
        slide,
        Inches(0.98),
        Inches(2.12),
        Inches(5.0),
        Inches(4.2),
        [
            "API 层：负责协议、鉴权入口、参数校验和响应封装",
            "Service 层：承接实体级业务能力，例如 Project / Series / Asset / Frame / VideoTask",
            "Workflow 层：串联跨步骤流程，例如资产生成、故事板生成、媒体生产",
            "Provider 层：统一模型、图像、视频、音频、导出、文本处理等外部能力接入",
            "Repository 层：对接数据库模型与领域对象映射，为持久化打地基",
        ],
        14,
        TEXT,
    )

    right_box = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(6.72), Inches(1.35), Inches(5.88), Inches(5.55))
    right_box.fill.solid()
    right_box.fill.fore_color.rgb = WHITE
    right_box.line.color.rgb = DIVIDER
    add_textbox(slide, Inches(6.98), Inches(1.6), Inches(5.2), Inches(0.45), ["直接收益"], 20, PRIMARY, True)
    add_textbox(
        slide,
        Inches(6.98),
        Inches(2.12),
        Inches(5.15),
        Inches(3.8),
        [
            "• 代码职责更清晰，后续新功能更容易定位改动点",
            "• 大幅降低单文件膨胀和“改一处牵全身”的连锁风险",
            "• 更适合多人并行开发、分模块测试和问题排查",
            "• 为数据库化、多租户化、平台化输出预留了清晰接口",
        ],
        15,
        TEXT,
    )
    add_metric(slide, Inches(7.0), Inches(5.35), Inches(1.6), Inches(1.0), "10", "服务模块", ACCENT)
    add_metric(slide, Inches(8.92), Inches(5.35), Inches(1.6), Inches(1.0), "4", "工作流模块", SECONDARY)
    add_metric(slide, Inches(10.84), Inches(5.35), Inches(1.6), Inches(1.0), "12", "仓储文件", ACCENT_2)

    # Slide 5
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_full_background(slide, BG)
    add_header_bar(slide, "核心改造二：产品链路与创作体验增强")

    add_card(
        slide,
        Inches(0.72),
        Inches(1.18),
        Inches(3.9),
        Inches(2.2),
        "系列化创作能力",
        [
            "• 引入 Series 作为分集项目的上层容器",
            "• 支持共享角色/场景/道具与系列级模型、提示词配置",
            "• 支持文件导入、分集预拆分和跨阶段资产复用",
        ],
        ACCENT,
    )
    add_card(
        slide,
        Inches(4.72),
        Inches(1.18),
        Inches(3.9),
        Inches(2.2),
        "提示词与创作辅助",
        [
            "• 增加多轮交互式提示词润色",
            "• 支持项目级与系列级 Prompt Config",
            "• 将最佳实践嵌入 Storyboard / Motion 等关键步骤",
        ],
        SECONDARY,
    )
    add_card(
        slide,
        Inches(8.72),
        Inches(1.18),
        Inches(3.9),
        Inches(2.2),
        "成片导出与交互完善",
        [
            "• 补齐真实导出接口与步骤 7-9 最小可用链路",
            "• 增加 TTS 失败反馈、导航重构、布局与透明视觉统一",
            "• 优化 VideoCreator、参数配置和全局可用性体验",
        ],
        ACCENT_2,
    )

    add_image(slide, IMAGES_DIR / "Script_example.jpg", Inches(0.82), Inches(3.75), 2.7, 2.1)
    add_image(slide, IMAGES_DIR / "StoryBoard_example.jpg", Inches(3.65), Inches(3.75), 2.7, 2.1)
    add_image(slide, IMAGES_DIR / "Motion_example.jpg", Inches(6.48), Inches(3.75), 2.7, 2.1)
    add_image(slide, IMAGES_DIR / "Assembly_example.jpg", Inches(9.31), Inches(3.75), 2.7, 2.1)
    add_textbox(slide, Inches(0.9), Inches(5.9), Inches(11.5), Inches(0.45), ["产品价值：让用户在单平台内完成“写、拆、画、动、拼、导”的完整生产闭环。"], 14, PRIMARY, True)

    # Slide 6
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_full_background(slide, BG)
    add_header_bar(slide, "核心改造三：数据库与配置基础设施升级")

    add_card(
        slide,
        Inches(0.72),
        Inches(1.25),
        Inches(4.0),
        Inches(4.9),
        "数据库能力演进",
        [
            "• 补齐 SQLAlchemy ORM 模型，覆盖 Project / Series / Character / Scene / Prop / Variant 等核心实体",
            "• 统一主键、审计字段、索引与 JSON/JSONB 结构定义",
            "• 当前已从 MySQL 初始化脚本转向 PostgreSQL 兼容 schema 与 session 配置",
        ],
        SECONDARY,
    )
    add_card(
        slide,
        Inches(4.92),
        Inches(1.25),
        Inches(3.55),
        Inches(4.9),
        "配置管理收口",
        [
            "• 新增 env_settings 作为 .env 统一读写入口",
            "• 同时兼容开发态和打包应用态配置路径",
            "• 系统启动时自动完成配置刷新、数据库初始化和输出目录准备",
        ],
        ACCENT,
    )
    add_card(
        slide,
        Inches(8.67),
        Inches(1.25),
        Inches(3.95),
        Inches(4.9),
        "测试与验证",
        [
            "• 新增仓储层 round-trip 测试，验证对象写入/读取/删除",
            "• 新增 PostgreSQL schema 与 DATABASE_URL 组装测试",
            "• 为后续数据库迁移和线上稳定性提供基本护栏",
        ],
        ACCENT_2,
    )

    # Slide 7
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_full_background(slide, BG)
    add_header_bar(slide, "阶段性成效与管理价值")

    add_metric(slide, Inches(0.9), Inches(1.4), Inches(2.2), Inches(1.15), "4", "主干关键提交", ACCENT)
    add_metric(slide, Inches(3.45), Inches(1.4), Inches(2.2), Inches(1.15), "82", "累计改动文件", SECONDARY)
    add_metric(slide, Inches(6.0), Inches(1.4), Inches(2.2), Inches(1.15), "62", "当前在制文件", ACCENT_2)
    add_metric(slide, Inches(8.55), Inches(1.4), Inches(2.2), Inches(1.15), "431+", "当前新增行", ACCENT)
    add_metric(slide, Inches(11.1), Inches(1.4), Inches(1.5), Inches(1.15), "816-", "当前删减行", PRIMARY)

    add_card(
        slide,
        Inches(0.72),
        Inches(3.0),
        Inches(3.9),
        Inches(2.7),
        "对研发团队的价值",
        [
            "• 降低系统复杂度，便于新成员接手",
            "• 提升模块边界清晰度，缩短排障路径",
            "• 支撑后续按模块排期与并行交付",
        ],
        ACCENT,
    )
    add_card(
        slide,
        Inches(4.72),
        Inches(3.0),
        Inches(3.9),
        Inches(2.7),
        "对产品推进的价值",
        [
            "• 支持系列化内容生产，增强平台复用性",
            "• 端到端链路更完整，适合对外演示与试用",
            "• 交互和导出能力补齐后，更接近可交付状态",
        ],
        SECONDARY,
    )
    add_card(
        slide,
        Inches(8.72),
        Inches(3.0),
        Inches(3.9),
        Inches(2.7),
        "对后续商业化/平台化的价值",
        [
            "• 多租户组织字段、工作区字段已在模型层预埋",
            "• 数据库结构化后，更利于权限、计费、协同能力接入",
            "• 为后续云化部署与长期运维打下基础",
        ],
        ACCENT_2,
    )

    # Slide 8
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_full_background(slide, BG)
    add_header_bar(slide, "下一阶段建议")
    add_card(
        slide,
        Inches(0.85),
        Inches(1.45),
        Inches(3.8),
        Inches(4.8),
        "1. 收口当前在制改造",
        [
            "• 完成 PostgreSQL schema、依赖与初始化流程的最终收敛",
            "• 跑通本地/测试环境验证，确保主链路不回退",
            "• 整理数据迁移说明，避免环境切换时出现断层",
        ],
        ACCENT,
    )
    add_card(
        slide,
        Inches(4.78),
        Inches(1.45),
        Inches(3.8),
        Inches(4.8),
        "2. 继续做平台级能力",
        [
            "• 深化组织、工作区、成员角色等多租户能力",
            "• 完善仓储层与接口层测试，提升发布信心",
            "• 补齐监控、日志与异常追踪，提高线上可运维性",
        ],
        SECONDARY,
    )
    add_card(
        slide,
        Inches(8.71),
        Inches(1.45),
        Inches(3.8),
        Inches(4.8),
        "3. 结合业务继续迭代",
        [
            "• 优化素材一致性、生成稳定性与成片质量",
            "• 针对样板客户/内部试用场景继续打磨易用性",
            "• 将技术改造沉淀为可复用的产品竞争力",
        ],
        ACCENT_2,
    )

    return prs


def main() -> None:
    prs = build_presentation()
    prs.save(OUTPUT_PATH)
    print(OUTPUT_PATH)


if __name__ == "__main__":
    main()
