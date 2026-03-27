from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Cm, Pt


# 中文字体统一使用更适合商务场景的常见字体。
FONT_FAMILY = "Microsoft YaHei"
TITLE_FONT_FAMILY = "Microsoft YaHei"

# 统一定义演讲稿的主色，保证整套页面风格一致。
NAVY = RGBColor(16, 34, 74)
BLUE = RGBColor(36, 99, 235)
CYAN = RGBColor(14, 165, 233)
TEAL = RGBColor(13, 148, 136)
GOLD = RGBColor(245, 158, 11)
SLATE = RGBColor(71, 85, 105)
LIGHT_BG = RGBColor(245, 247, 250)
WHITE = RGBColor(255, 255, 255)
DARK = RGBColor(15, 23, 42)
MUTED = RGBColor(100, 116, 139)
LIGHT_BORDER = RGBColor(226, 232, 240)


def set_run_style(run, *, size=18, bold=False, color=DARK, font=FONT_FAMILY):
    """统一设置文本 run 的样式，避免单页单独调整时风格漂移。"""
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def add_textbox(slide, left, top, width, height, text, *, size=18, bold=False,
                color=DARK, align=PP_ALIGN.LEFT, font=FONT_FAMILY):
    """创建文本框并写入一段文本。"""
    box = slide.shapes.add_textbox(left, top, width, height)
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.vertical_anchor = MSO_ANCHOR.TOP
    paragraph = frame.paragraphs[0]
    paragraph.alignment = align
    run = paragraph.add_run()
    run.text = text
    set_run_style(run, size=size, bold=bold, color=color, font=font)
    return box


def add_bullet_list(slide, left, top, width, height, items, *, font_size=20,
                    color=DARK, bullet_color=BLUE, level_indents=None):
    """创建分层要点列表，适合汇报类页面。"""
    box = slide.shapes.add_textbox(left, top, width, height)
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    if level_indents is None:
        level_indents = {0: Pt(0), 1: Pt(18)}

    for idx, item in enumerate(items):
        text = item["text"]
        level = item.get("level", 0)
        paragraph = frame.paragraphs[0] if idx == 0 else frame.add_paragraph()
        paragraph.text = text
        paragraph.level = level
        paragraph.alignment = PP_ALIGN.LEFT
        paragraph.space_after = Pt(8)
        paragraph.space_before = Pt(0)
        paragraph.line_spacing = 1.15
        paragraph.bullet = True
        paragraph.left_margin = level_indents.get(level, Pt(18))
        paragraph.hanging = Pt(0)
        for run in paragraph.runs:
            set_run_style(run, size=font_size - level * 2, color=color)
        if paragraph.runs:
            paragraph.runs[0].font.color.rgb = bullet_color if level == 0 else color
    return box


def add_card(slide, left, top, width, height, title, body_lines, *,
             accent=BLUE, title_size=20, body_size=15):
    """创建信息卡片，用于展示场景、能力和价值。"""
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = WHITE
    shape.line.color.rgb = LIGHT_BORDER
    shape.line.width = Pt(1.1)

    accent_bar = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, left, top, Cm(0.28), height)
    accent_bar.fill.solid()
    accent_bar.fill.fore_color.rgb = accent
    accent_bar.line.fill.background()

    add_textbox(slide, left + Cm(0.7), top + Cm(0.45), width - Cm(1.1), Cm(0.8), title,
                size=title_size, bold=True, color=NAVY)

    items = [{"text": line, "level": 0} for line in body_lines]
    add_bullet_list(slide, left + Cm(0.65), top + Cm(1.35), width - Cm(1.0), height - Cm(1.6),
                    items, font_size=body_size, color=SLATE, bullet_color=accent)
    return shape


def add_header(slide, title, subtitle=None, section_label=None):
    """为内容页绘制统一页眉。"""
    band = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, Cm(33.867), Cm(1.05))
    band.fill.solid()
    band.fill.fore_color.rgb = NAVY
    band.line.fill.background()

    if section_label:
        pill = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Cm(1.2), Cm(1.3), Cm(2.4), Cm(0.7))
        pill.fill.solid()
        pill.fill.fore_color.rgb = LIGHT_BG
        pill.line.fill.background()
        add_textbox(slide, Cm(1.35), Cm(1.42), Cm(2.1), Cm(0.45), section_label,
                    size=10, bold=True, color=BLUE, align=PP_ALIGN.CENTER)

    add_textbox(slide, Cm(1.2), Cm(2.0), Cm(20), Cm(1.2), title, size=26, bold=True, color=NAVY, font=TITLE_FONT_FAMILY)
    if subtitle:
        add_textbox(slide, Cm(1.2), Cm(3.05), Cm(24), Cm(0.8), subtitle, size=12, color=MUTED)


def add_footer(slide, page_no):
    """统一添加页脚和页码，便于商务场合投屏。"""
    line = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Cm(1.2), Cm(18.45), Cm(31.2), Pt(1.5))
    line.fill.solid()
    line.fill.fore_color.rgb = LIGHT_BORDER
    line.line.fill.background()
    add_textbox(slide, Cm(1.2), Cm(18.6), Cm(8), Cm(0.4), "LumenX Studio | AI赋能教育与组织效率", size=9, color=MUTED)
    add_textbox(slide, Cm(31.0), Cm(18.5), Cm(1.2), Cm(0.5), str(page_no), size=10, color=MUTED, align=PP_ALIGN.RIGHT)


def build_cover(prs):
    """封面页用于建立正式、可信的商务第一印象。"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = WHITE
    bg.line.fill.background()

    left_panel = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, Cm(11.5), prs.slide_height)
    left_panel.fill.solid()
    left_panel.fill.fore_color.rgb = NAVY
    left_panel.line.fill.background()

    accent = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Cm(11.5), 0, Cm(0.25), prs.slide_height)
    accent.fill.solid()
    accent.fill.fore_color.rgb = CYAN
    accent.line.fill.background()

    add_textbox(slide, Cm(1.3), Cm(3.2), Cm(8.5), Cm(0.9), "商务汇报", size=16, bold=True, color=WHITE)
    add_textbox(slide, Cm(1.3), Cm(4.4), Cm(17), Cm(2.2), "人工智能在教育与组织场景中的应用", size=28, bold=True, color=WHITE, font=TITLE_FONT_FAMILY)
    add_textbox(slide, Cm(1.3), Cm(7.0), Cm(8.8), Cm(1.0), "从教学提效到流程自动化的落地路径", size=15, color=RGBColor(214, 224, 255))

    quote = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Cm(13.1), Cm(4.2), Cm(17.4), Cm(4.2))
    quote.fill.solid()
    quote.fill.fore_color.rgb = LIGHT_BG
    quote.line.fill.background()
    add_textbox(slide, Cm(14.0), Cm(5.0), Cm(15.8), Cm(2.2),
                "本次演讲聚焦两个问题：\n1. AI 如何重塑课前、课中、课后教学流程\n2. AI 如何在数据分析、审批、内容生成中形成业务价值",
                size=18, color=NAVY)

    tag1 = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Cm(13.2), Cm(10.0), Cm(5.2), Cm(1.0))
    tag1.fill.solid()
    tag1.fill.fore_color.rgb = RGBColor(226, 239, 255)
    tag1.line.fill.background()
    add_textbox(slide, Cm(13.55), Cm(10.25), Cm(4.6), Cm(0.5), "教育数字化升级", size=13, bold=True, color=BLUE, align=PP_ALIGN.CENTER)

    tag2 = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Cm(19.0), Cm(10.0), Cm(4.8), Cm(1.0))
    tag2.fill.solid()
    tag2.fill.fore_color.rgb = RGBColor(223, 250, 243)
    tag2.line.fill.background()
    add_textbox(slide, Cm(19.35), Cm(10.25), Cm(4.1), Cm(0.5), "运营效率提升", size=13, bold=True, color=TEAL, align=PP_ALIGN.CENTER)

    tag3 = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Cm(24.3), Cm(10.0), Cm(5.0), Cm(1.0))
    tag3.fill.solid()
    tag3.fill.fore_color.rgb = RGBColor(255, 243, 224)
    tag3.line.fill.background()
    add_textbox(slide, Cm(24.65), Cm(10.25), Cm(4.3), Cm(0.5), "可复制的应用场景", size=13, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
    return slide


def build_agenda(prs):
    """目录页帮助商务听众快速建立演讲预期。"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "汇报目录", "围绕教育主场景与通用AI能力两条主线展开", "Agenda")

    agenda_items = [
        ("01", "为什么现在要关注 AI+教育", "从提效工具走向教学与管理协同平台"),
        ("02", "课前、课中、课后的应用全景", "把教学过程拆解成可量化、可优化的节点"),
        ("03", "三个通用 AI 演示方向", "数据统计、流程审批、视频生成"),
        ("04", "落地建议与合作价值", "从试点到规模化推广的实施路径"),
    ]

    top = 4.3
    for idx, (num, title, desc) in enumerate(agenda_items):
        y = Cm(top + idx * 3.1)
        card = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Cm(1.4), y, Cm(30.6), Cm(2.3))
        card.fill.solid()
        card.fill.fore_color.rgb = WHITE
        card.line.color.rgb = LIGHT_BORDER
        num_box = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Cm(1.9), y + Cm(0.45), Cm(2.0), Cm(1.35))
        num_box.fill.solid()
        num_box.fill.fore_color.rgb = [BLUE, TEAL, CYAN, GOLD][idx]
        num_box.line.fill.background()
        add_textbox(slide, Cm(2.2), y + Cm(0.78), Cm(1.4), Cm(0.5), num, size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_textbox(slide, Cm(4.7), y + Cm(0.42), Cm(12), Cm(0.7), title, size=19, bold=True, color=NAVY)
        add_textbox(slide, Cm(4.7), y + Cm(1.15), Cm(18), Cm(0.6), desc, size=11, color=MUTED)

    add_footer(slide, 2)
    return slide


def build_why_now(prs):
    """强调业务价值，帮助商务受众理解演讲的现实意义。"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "为什么现在是 AI+教育 的关键窗口期", "AI 不只是展示型功能，而是贯穿教学与管理流程的生产力工具", "Value")

    add_card(
        slide, Cm(1.2), Cm(4.3), Cm(9.7), Cm(5.4), "教学内容生产提效",
        ["备课从人工整理，转向素材自动生成与快速迭代", "视频、图片、实验演示内容可按课程即时定制", "缩短教师从想法到课堂落地的准备周期"],
        accent=BLUE
    )
    add_card(
        slide, Cm(11.5), Cm(4.3), Cm(9.7), Cm(5.4), "课堂过程可互动、可分析",
        ["课堂提问与互动内容可以即时生成与推送", "课堂行为数据与互动结果可以自动汇总", "让“课堂效果”从经验判断变成数据支持"],
        accent=TEAL
    )
    add_card(
        slide, Cm(21.8), Cm(4.3), Cm(9.0), Cm(5.4), "校内管理与运营自动化",
        ["数据统计、审批流、内容生产具备可复制性", "减少重复劳动，释放组织管理效率", "形成平台化能力，而非单点工具"],
        accent=GOLD
    )

    band = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Cm(1.2), Cm(11.1), Cm(29.6), Cm(3.4))
    band.fill.solid()
    band.fill.fore_color.rgb = LIGHT_BG
    band.line.fill.background()
    add_textbox(slide, Cm(1.8), Cm(11.75), Cm(28.5), Cm(1.0), "核心结论：AI 最适合优先切入“高频、重复、标准化程度高”的教育与管理环节，再逐步向个性化教学支持延伸。", size=20, bold=True, color=NAVY)
    add_textbox(slide, Cm(1.8), Cm(13.1), Cm(28.5), Cm(0.7), "这意味着项目落地时应先做试点闭环，再做规模复制。", size=12, color=MUTED)

    add_footer(slide, 3)
    return slide


def build_education_overview(prs):
    """总览页用于把脑图结构转成商务表达。"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "AI+教育应用全景", "以教学全生命周期为主线：课前、课中、课后", "Education")

    center_left = Cm(12.9)
    center_top = Cm(7.1)
    core = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, center_left, center_top, Cm(7.8), Cm(2.0))
    core.fill.solid()
    core.fill.fore_color.rgb = NAVY
    core.line.fill.background()
    add_textbox(slide, center_left + Cm(0.4), center_top + Cm(0.55), Cm(7.0), Cm(0.6), "人工智能 + 教育", size=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    add_card(slide, Cm(1.2), Cm(4.6), Cm(8.6), Cm(4.8), "课前：内容准备", ["备课内容输出", "实验/情景/知识点素材生成", "教师准备效率提升"], accent=BLUE)
    add_card(slide, Cm(1.2), Cm(10.4), Cm(8.6), Cm(4.6), "课后：学习闭环", ["作业布置", "主观题与客观题自动批改", "学习结果快速反馈"], accent=CYAN)
    add_card(slide, Cm(23.0), Cm(4.6), Cm(8.6), Cm(4.8), "课中：互动与评价", ["课堂互动即时触发", "课堂行为分析", "课堂效果输出"], accent=TEAL)
    add_card(slide, Cm(23.0), Cm(10.4), Cm(8.6), Cm(4.6), "延展能力", ["AI体育", "管理系统", "校园治理与平台协同"], accent=GOLD)

    add_footer(slide, 4)
    return slide


def build_pre_class(prs):
    """课前场景页突出内容生产价值。"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "课前：让备课从“找素材”升级为“生产内容”", "围绕脑图中的备课内容输出展开", "Pre-Class")

    add_card(slide, Cm(1.2), Cm(4.3), Cm(9.5), Cm(9.6), "场景一：物理实验游戏", [
        "AI 根据知识点自动生成实验演示与互动玩法",
        "适合把抽象原理转化为可视化、可参与的体验",
        "价值：提高理解效率，降低教师素材制作成本",
    ], accent=BLUE, body_size=15)

    add_card(slide, Cm(11.6), Cm(4.3), Cm(9.5), Cm(9.6), "场景二：语文情景式教学视频", [
        "围绕古诗、课文或文学场景生成画面与短视频",
        "把文字意境转成情境化教学内容",
        "价值：增强课堂代入感，提升学生注意力",
    ], accent=TEAL, body_size=15)

    add_card(slide, Cm(22.0), Cm(4.3), Cm(9.0), Cm(9.6), "场景三：计算机组成图像/视频", [
        "用 AI 生成 CPU、ROM、总线等模块示意图与讲解视频",
        "适用于职业教育、信息技术课程等结构性知识教学",
        "价值：统一教学素材标准，提升复杂知识表达效率",
    ], accent=GOLD, body_size=15)

    add_footer(slide, 5)
    return slide


def build_in_class(prs):
    """课中场景页突出互动与评估闭环。"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "课中：把课堂互动与教学评价做成实时闭环", "对应脑图中的课堂互动与课堂评价", "In-Class")

    left = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Cm(1.2), Cm(4.6), Cm(14.7), Cm(9.2))
    left.fill.solid()
    left.fill.fore_color.rgb = WHITE
    left.line.color.rgb = LIGHT_BORDER
    add_textbox(slide, Cm(1.9), Cm(5.2), Cm(6), Cm(0.8), "课堂互动示例", size=22, bold=True, color=NAVY)
    add_textbox(slide, Cm(1.9), Cm(6.2), Cm(12.6), Cm(2.4),
                "老师讲完“计算机组成原理”后，即时发起问题：\n“ROM 属于计算机的哪个模块？”\n学生端同步弹出互动卡片完成答题。",
                size=17, color=SLATE)

    teacher = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Cm(2.2), Cm(9.0), Cm(5.2), Cm(3.4))
    teacher.fill.solid()
    teacher.fill.fore_color.rgb = RGBColor(231, 240, 255)
    teacher.line.fill.background()
    add_textbox(slide, Cm(2.65), Cm(9.55), Cm(4.3), Cm(0.8), "教师端\n发布互动问题", size=19, bold=True, color=BLUE, align=PP_ALIGN.CENTER)

    student = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Cm(9.0), Cm(9.0), Cm(5.2), Cm(3.4))
    student.fill.solid()
    student.fill.fore_color.rgb = RGBColor(223, 250, 243)
    student.line.fill.background()
    add_textbox(slide, Cm(9.45), Cm(9.55), Cm(4.3), Cm(0.8), "学生端\n接收并回答", size=19, bold=True, color=TEAL, align=PP_ALIGN.CENTER)

    arrow = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.CHEVRON, Cm(7.6), Cm(10.1), Cm(1.0), Cm(1.2))
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = CYAN
    arrow.line.fill.background()

    add_card(slide, Cm(17.0), Cm(4.6), Cm(14.0), Cm(9.2), "课堂评价输出", [
        "AI 分析教师课堂行为、讲授节奏和学生互动反馈",
        "把课堂参与度、响应情况、知识点理解情况形成结果输出",
        "帮助管理层、教研组和教师本人做复盘优化",
        "从“感觉课堂不错”升级为“课堂效果可观测、可追踪”",
    ], accent=TEAL, body_size=16)

    add_footer(slide, 6)
    return slide


def build_post_class(prs):
    """课后页补足教学闭环，同时保留脑图中的其他扩展能力。"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "课后与延展：形成学习闭环，并向校园管理外溢", "覆盖作业、批改、AI体育与管理系统", "Post-Class")

    add_card(slide, Cm(1.2), Cm(4.3), Cm(9.6), Cm(8.8), "作业布置", [
        "根据课程内容自动生成作业建议与题目组合",
        "支持因材施教的分层作业设计",
        "让教师从重复性出题中解放出来",
    ], accent=BLUE)

    add_card(slide, Cm(11.8), Cm(4.3), Cm(9.6), Cm(8.8), "AI 自动批改", [
        "覆盖客观题快速批改",
        "结合规则与模型能力辅助主观题评价",
        "缩短反馈周期，提升学生改进效率",
    ], accent=CYAN)

    add_card(slide, Cm(22.4), Cm(4.3), Cm(8.7), Cm(4.1), "AI 体育", [
        "动作识别、训练辅助、过程评价",
        "拓展至体育教学与校园运动场景",
    ], accent=TEAL, body_size=14)

    add_card(slide, Cm(22.4), Cm(9.0), Cm(8.7), Cm(4.1), "管理系统", [
        "与校务、数据、流程系统协同",
        "让 AI 能力从课堂延伸到学校运营层",
    ], accent=GOLD, body_size=14)

    add_footer(slide, 7)
    return slide


def build_general_ai(prs):
    """通用 AI 应用页用于连接教育场景与管理场景。"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "人工智能应用演示方向", "将教育能力延展到组织管理与内容生产场景", "Demo")

    demo_data = [
        ("演示 1", "数据统计", "例如招聘数据统计", "自动汇总多维数据，快速输出分析结论", BLUE),
        ("演示 2", "流程审批", "例如批量处理请假审批", "把重复性审批流交给 AI 辅助处理", TEAL),
        ("演示 3", "视频生成", "例如教学视频生成", "将内容生产从文本转向多媒体自动化生成", GOLD),
    ]

    for i, (num, title, example, desc, color) in enumerate(demo_data):
        x = Cm(1.2 + i * 10.4)
        box = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x, Cm(5.1), Cm(9.4), Cm(8.8))
        box.fill.solid()
        box.fill.fore_color.rgb = WHITE
        box.line.color.rgb = LIGHT_BORDER

        top_band = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, x, Cm(5.1), Cm(9.4), Cm(1.1))
        top_band.fill.solid()
        top_band.fill.fore_color.rgb = color
        top_band.line.fill.background()

        add_textbox(slide, x + Cm(0.45), Cm(5.43), Cm(2.0), Cm(0.4), num, size=12, bold=True, color=WHITE)
        add_textbox(slide, x + Cm(0.55), Cm(6.65), Cm(6.5), Cm(0.8), title, size=24, bold=True, color=NAVY)
        add_textbox(slide, x + Cm(0.55), Cm(7.75), Cm(7.8), Cm(0.8), example, size=14, color=MUTED)
        add_textbox(slide, x + Cm(0.55), Cm(9.2), Cm(8.0), Cm(2.0), desc, size=17, color=SLATE)
        add_textbox(slide, x + Cm(0.55), Cm(12.0), Cm(7.9), Cm(1.0), "适合作为现场演示或试点切入口", size=12, bold=True, color=color)

    add_footer(slide, 8)
    return slide


def build_roadmap(prs):
    """给商务场合增加落地建议，让 PPT 更完整。"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "建议的落地路径", "从试点验证到规模推广，降低实施风险", "Roadmap")

    steps = [
        ("阶段 1", "试点场景验证", "选择 1-2 个最有感知度的场景，例如备课视频生成、课堂互动、审批自动化"),
        ("阶段 2", "形成数据闭环", "沉淀使用数据、课堂效果、效率收益，为后续复制提供量化依据"),
        ("阶段 3", "平台化扩展", "把内容生成、互动、评价、管理能力统一到平台中，服务更多校区/部门"),
    ]

    for idx, (stage, title, desc) in enumerate(steps):
        y = Cm(5.2 + idx * 3.35)
        circle = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Cm(1.4), y, Cm(2.0), Cm(2.0))
        circle.fill.solid()
        circle.fill.fore_color.rgb = [BLUE, TEAL, GOLD][idx]
        circle.line.fill.background()
        add_textbox(slide, Cm(1.7), y + Cm(0.62), Cm(1.4), Cm(0.7), str(idx + 1), size=20, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

        panel = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Cm(4.0), y - Cm(0.1), Cm(26.8), Cm(2.2))
        panel.fill.solid()
        panel.fill.fore_color.rgb = WHITE
        panel.line.color.rgb = LIGHT_BORDER
        add_textbox(slide, Cm(4.6), y + Cm(0.18), Cm(3.0), Cm(0.5), stage, size=12, bold=True, color=[BLUE, TEAL, GOLD][idx])
        add_textbox(slide, Cm(4.6), y + Cm(0.72), Cm(8.0), Cm(0.6), title, size=19, bold=True, color=NAVY)
        add_textbox(slide, Cm(13.4), y + Cm(0.68), Cm(16.4), Cm(0.8), desc, size=13, color=SLATE)

    add_footer(slide, 9)
    return slide


def build_closing(prs, thumbnail_path: Path):
    """结尾页回收到合作与价值表达。"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = NAVY
    bg.line.fill.background()

    add_textbox(slide, Cm(1.5), Cm(3.2), Cm(12), Cm(0.9), "总结", size=16, bold=True, color=RGBColor(170, 190, 255))
    add_textbox(slide, Cm(1.5), Cm(4.5), Cm(16), Cm(2.0), "AI 不是替代教学，\n而是放大优质教学与高效管理。", size=28, bold=True, color=WHITE, font=TITLE_FONT_FAMILY)
    add_textbox(slide, Cm(1.5), Cm(7.6), Cm(14.5), Cm(2.0), "建议以“可见效果、可量化收益、可复制扩展”为标准，\n选择试点场景快速启动。", size=16, color=RGBColor(218, 226, 255))

    if thumbnail_path.exists():
        slide.shapes.add_picture(str(thumbnail_path), Cm(19.4), Cm(3.8), width=Cm(11.6))
        frame = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Cm(19.0), Cm(3.4), Cm(12.4), Cm(8.5))
        frame.fill.background()
        frame.line.color.rgb = RGBColor(148, 163, 184)
        frame.line.width = Pt(1.0)
        slide.shapes._spTree.remove(frame._element)
        slide.shapes._spTree.append(frame._element)

    add_textbox(slide, Cm(1.5), Cm(15.9), Cm(8), Cm(0.6), "谢谢", size=18, bold=True, color=WHITE)
    add_textbox(slide, Cm(1.5), Cm(16.7), Cm(12), Cm(0.5), "期待进一步交流试点方案与合作路径", size=11, color=RGBColor(189, 200, 230))
    return slide


def main():
    """生成 PPT 主流程。"""
    output_dir = Path("/Users/will/Documents/jishu/code/lumenx/output/ppt")
    output_dir.mkdir(parents=True, exist_ok=True)
    output_path = output_dir / "AI教育与应用商务演讲稿.pptx"
    thumbnail_path = output_dir / "Thumbnails" / "thumbnail.png"

    prs = Presentation()
    prs.slide_width = Cm(33.867)
    prs.slide_height = Cm(19.05)

    build_cover(prs)
    build_agenda(prs)
    build_why_now(prs)
    build_education_overview(prs)
    build_pre_class(prs)
    build_in_class(prs)
    build_post_class(prs)
    build_general_ai(prs)
    build_roadmap(prs)
    build_closing(prs, thumbnail_path)

    prs.save(output_path)
    print(output_path)


if __name__ == "__main__":
    main()
